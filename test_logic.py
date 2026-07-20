import json
import io
import gc
import importlib
import logging
import tempfile
import threading
import unittest
from contextlib import redirect_stdout
from datetime import datetime, timedelta
from http.server import BaseHTTPRequestHandler, ThreadingHTTPServer
from pathlib import Path
from unittest.mock import patch

import api
from backend.routes.api import JsApi, _session_path
from backend.services import document_service as documents
from backend.services import codex_limits
from backend.services import local_llm
from backend.services import llm_service as llm
from backend.services import json_tools
from backend.services import translation_service
from backend.services import workflow_service as workflows
from config import settings
import backend.window_main as wm
from pet.state_machine import PetState, PetStateMachine
from plugins import builtin as builtin_plugins
import scheduler
import spritecat
import worker


class SchedulerTests(unittest.TestCase):
    def test_invalid_json_is_reported(self):
        with tempfile.TemporaryDirectory() as directory:
            path = Path(directory) / 'schedule.json'
            path.write_text('{bad json', encoding='utf-8')
            result = scheduler.Scheduler(path)
            self.assertTrue(result.errors)
            self.assertEqual(result.items, [])

    def test_daily_lead_and_ontime_fire_once(self):
        with tempfile.TemporaryDirectory() as directory:
            path = Path(directory) / 'schedule.json'
            now = datetime(2026, 7, 17, 9, 0)
            path.write_text(json.dumps([{
                'id': 'daily-test', 'title': 'Daily Test', 'type': 'daily',
                'time': '09:02', 'lead_min': 1, 'enabled': True,
            }]), encoding='utf-8')
            result = scheduler.Scheduler(path)

            self.assertEqual([(item['id'], kind) for item, kind in result.tick(now + timedelta(minutes=1))],
                             [('daily-test', 'lead')])
            self.assertEqual(result.tick(now + timedelta(minutes=1)), [])
            self.assertEqual([(item['id'], kind) for item, kind in result.tick(now + timedelta(minutes=2))],
                             [('daily-test', 'ontime')])

    def test_weekly_hourly_and_cross_midnight_lead(self):
        with tempfile.TemporaryDirectory() as directory:
            path = Path(directory) / 'schedule.json'
            path.write_text(json.dumps([
                {'id': 'weekly', 'title': 'Weekly', 'type': 'weekly', 'day': 'MO',
                 'time': '09:00', 'lead_min': 0, 'enabled': True},
                {'id': 'hourly', 'title': 'Hourly', 'type': 'hourly', 'minute': 30,
                 'lead_min': 0, 'enabled': True},
                {'id': 'midnight', 'title': 'Midnight', 'type': 'daily', 'time': '00:01',
                 'lead_min': 2, 'enabled': True},
            ]), encoding='utf-8')
            result = scheduler.Scheduler(path)
            self.assertEqual([(item['id'], kind) for item, kind in result.tick(datetime(2026, 7, 20, 9, 0))],
                             [('weekly', 'ontime')])
            self.assertEqual([(item['id'], kind) for item, kind in result.tick(datetime(2026, 7, 20, 13, 30))],
                             [('hourly', 'ontime')])
            self.assertEqual([(item['id'], kind) for item, kind in result.tick(datetime(2026, 7, 20, 23, 59))],
                             [('midnight', 'lead')])
            self.assertEqual([(item['id'], kind) for item, kind in result.tick(datetime(2026, 7, 21, 0, 1))],
                             [('midnight', 'ontime')])

    def test_weekly_non_matching_day_does_not_fire(self):
        with tempfile.TemporaryDirectory() as directory:
            path = Path(directory) / 'schedule.json'
            path.write_text(json.dumps([{
                'id': 'weekly', 'title': 'Weekly', 'type': 'weekly', 'day': 'MO',
                'time': '09:00', 'lead_min': 0, 'enabled': True,
            }]), encoding='utf-8')
            result = scheduler.Scheduler(path)
            self.assertEqual(result.tick(datetime(2026, 7, 21, 9, 0)), [])  # Tuesday

    def test_disabled_and_deleted_schedule_never_fires(self):
        with tempfile.TemporaryDirectory() as directory:
            path = Path(directory) / 'schedule.json'
            path.write_text(json.dumps([{
                'id': 'disabled', 'title': 'Disabled', 'type': 'daily', 'time': '09:00',
                'lead_min': 0, 'enabled': False,
            }, {
                'id': 'delete-me', 'title': 'Delete me', 'type': 'daily', 'time': '09:00',
                'lead_min': 0, 'enabled': True,
            }]), encoding='utf-8')
            result = scheduler.Scheduler(path)
            result.delete('delete-me')
            self.assertEqual(result.tick(datetime(2026, 7, 20, 9, 0)), [])


class PetStateTests(unittest.TestCase):
    def test_chat_lifecycle_transitions_are_deterministic(self):
        state = PetStateMachine()
        for target in (PetState.LISTENING, PetState.THINKING, PetState.STREAMING, PetState.SUCCESS):
            self.assertTrue(state.transition(target))
        self.assertEqual(state.current, PetState.SUCCESS)
        self.assertTrue(state.transition(PetState.IDLE))

    def test_invalid_transition_keeps_existing_state(self):
        state = PetStateMachine()
        self.assertFalse(state.transition(PetState.STREAMING))
        self.assertEqual(state.current, PetState.IDLE)

    def test_idle_sleep_does_not_depend_on_optional_usage_monitor(self):
        # cat.py creates its runtime log handler at import time.  Keep this
        # deterministic logic test independent of a running desktop app's log.
        with patch('logging.handlers.RotatingFileHandler', return_value=logging.NullHandler()):
            cat_module = importlib.import_module('cat')
        pet = object.__new__(cat_module.ClaudeCat)
        pet.sleep_frames = [object()]
        pet._sleep_min = 1
        pet._last_interact = 0
        pet._effective_usage = lambda: None
        self.assertTrue(pet._should_sleep())

    def test_restoring_full_size_clamps_a_docked_pet_on_screen(self):
        with patch('logging.handlers.RotatingFileHandler', return_value=logging.NullHandler()):
            cat_module = importlib.import_module('cat')
        pet = object.__new__(cat_module.ClaudeCat)
        pet.root = type('Root', (), {
            'winfo_screenwidth': lambda _self: 1366,
            'winfo_screenheight': lambda _self: 768,
        })()
        pet.w = pet.h = 128
        self.assertEqual(pet._clamp_pet_position(1334, 689), (1238, 640))

    def test_usage_badge_is_hidden_when_both_limits_are_off(self):
        with patch('logging.handlers.RotatingFileHandler', return_value=logging.NullHandler()):
            cat_module = importlib.import_module('cat')

        class Value:
            def __init__(self, value):
                self.value = value

            def get(self):
                return self.value

        class BadgeWindow:
            def __init__(self):
                self.withdrawn = False

            def withdraw(self):
                self.withdrawn = True

        pet = object.__new__(cat_module.ClaudeCat)
        pet.show_pct = Value(True)
        pet.monitor_enabled = Value(False)
        pet.codex_limits_enabled = Value(False)
        pet.badge_win = BadgeWindow()

        pet._update_badge()

        self.assertTrue(pet.badge_win.withdrawn)
        pet.codex_limits_enabled.value = True
        self.assertTrue(pet._usage_badge_enabled())

    def test_clear_log_removes_rotated_files_and_logging_continues(self):
        with patch('logging.handlers.RotatingFileHandler', return_value=logging.NullHandler()):
            cat_module = importlib.import_module('cat')
        with tempfile.TemporaryDirectory() as directory:
            log_file = Path(directory) / 'claudecat.log'
            test_logger = logging.getLogger(f'claudecat-clear-test-{id(self)}')
            test_logger.setLevel(logging.INFO)
            test_logger.propagate = False
            handler = logging.handlers.RotatingFileHandler(
                log_file, maxBytes=1024, backupCount=2, encoding='utf-8')
            test_logger.addHandler(handler)
            try:
                test_logger.info('before clear')
                Path(f'{log_file}.1').write_text('older log', encoding='utf-8')

                cat_module._clear_log_files(test_logger, log_file)

                self.assertEqual(log_file.read_text(encoding='utf-8'), '')
                self.assertFalse(Path(f'{log_file}.1').exists())
                test_logger.info('after clear')
                self.assertIn('after clear', log_file.read_text(encoding='utf-8'))
            finally:
                test_logger.removeHandler(handler)
                handler.close()


class PluginTests(unittest.TestCase):
    def test_builtin_plugins_are_fixed_allowlisted_actions(self):
        actions = builtin_plugins.actions()
        self.assertEqual([action.action_id for action in actions], ['quick_question', 'documents'])
        self.assertTrue(all(action.label for action in actions))


class SpriteStateTests(unittest.TestCase):
    def test_named_state_assets_do_not_leak_into_run_cycle(self):
        with tempfile.TemporaryDirectory() as directory:
            root = Path(directory)
            for name in ('bluecat_run_01.png', 'bluecat_idle_01.png',
                         'bluecat_listening_01.png', 'bluecat_thinking_01.png',
                         'bluecat_success_01.png', 'bluecat_error_01.png'):
                (root / name).touch()
            actions = spritecat._action_paths(root)
            self.assertEqual([path.name for path in actions['run']], ['bluecat_run_01.png'])
            self.assertEqual(len(actions['listening']), 1)
            self.assertEqual(len(actions['thinking']), 1)
            self.assertEqual(len(actions['success']), 1)


class ApiTests(unittest.TestCase):
    def test_read_access_token_uses_configured_credentials_path(self):
        with tempfile.TemporaryDirectory() as directory:
            credentials = Path(directory) / '.credentials.json'
            credentials.write_text(json.dumps({
                'claudeAiOauth': {'accessToken': 'test-token'},
            }), encoding='utf-8')
            with patch.object(api, 'CLAUDE_CREDENTIALS', credentials):
                self.assertEqual(api.read_access_token(), 'test-token')

    def test_session_path_rejects_traversal(self):
        self.assertIsNone(_session_path('../config'))
        self.assertIsNotNone(_session_path('550e8400-e29b-41d4-a716-446655440000'))

    def test_claude_usage_api_is_hard_disabled_before_http_request(self):
        api.set_usage_api_enabled(False)
        try:
            self.assertFalse(api.USAGE_API_ENABLED)
            with patch.object(api.requests, 'get') as request:
                result = api.fetch_usage()
            self.assertEqual(result['error'], 'Claude usage monitor disabled')
            request.assert_not_called()
        finally:
            api.set_usage_api_enabled(False)

    def test_claude_usage_api_toggle_changes_only_the_optional_monitor(self):
        api.set_usage_api_enabled(True)
        try:
            self.assertTrue(api.USAGE_API_ENABLED)
        finally:
            api.set_usage_api_enabled(False)

    def test_clearing_history_discards_an_inflight_response(self):
        wm.clear_history()
        with patch.object(llm, 'chat', side_effect=lambda *_args, **_kwargs: (
                wm.clear_history() or {'content': 'late answer'})):
            result = JsApi().send_message('hello')
        self.assertTrue(result['discarded'])
        _generation, history = wm.history_snapshot()
        self.assertEqual(history, [])

    def test_chat_uses_the_model_selected_by_the_user(self):
        wm._current_model = 'user-selected-model'
        with patch.object(wm, 'history_snapshot', return_value=(1, [])), \
             patch.object(wm, 'append_history_if_current', return_value=None), \
             patch.object(llm, 'chat', return_value={'content': 'ok'}) as chat:
            JsApi().send_message('hello')
        self.assertEqual(chat.call_args.kwargs['model'], 'user-selected-model')

    def test_every_assistant_tab_docks_the_pet_window(self):
        wm._open_evt.clear()
        try:
            with patch.object(wm, '_window', None), \
                 patch.object(wm, '_on_chat_open') as on_open:
                wm.request_open('schedule')
            on_open.assert_called_once()
        finally:
            wm._open_evt.clear()

    def test_tool_window_is_explicitly_shown_in_the_taskbar(self):
        native = type('NativeWindow', (), {'ShowInTaskbar': False})()
        window = type('ToolWindow', (), {'native': native})()

        wm._show_in_taskbar(window)

        self.assertTrue(native.ShowInTaskbar)


class LlmTests(unittest.TestCase):
    def test_model_list_deduplicates_primary_and_fallbacks(self):
        with patch.object(llm, '_config', {
            'model': 'primary',
            'fallback_models': ['fallback', 'primary', 'fallback'],
        }):
            self.assertEqual(llm.list_models(), ['primary', 'fallback'])

    def test_selecting_model_preserves_the_other_available_models(self):
        with tempfile.TemporaryDirectory() as directory:
            config_path = Path(directory) / 'config.json'
            with patch.object(settings, 'CONFIG_FILE', config_path), \
                 patch.object(settings, 'config', {}), \
                 patch.object(llm, '_config_file', config_path), \
                 patch.object(llm, '_config', {
                     'model': 'primary',
                     'fallback_models': ['fallback'],
                 }):
                llm.save_config_model('fallback')
                saved = json.loads(config_path.read_text(encoding='utf-8'))['llm']
                models = llm.list_models()
        self.assertEqual(saved['model'], 'fallback')
        self.assertEqual(saved['fallback_models'], ['primary'])
        self.assertEqual(models, ['fallback', 'primary'])

    def test_context_overflow_detection(self):
        self.assertTrue(llm._looks_like_context_overflow('Maximum context length exceeded'))
        self.assertFalse(llm._looks_like_context_overflow('Authentication failed'))

    def test_file_size_limit_has_a_safe_default(self):
        with patch.object(llm, '_config', {}):
            self.assertEqual(llm.max_file_bytes(), 10 * 1024 * 1024)

    def test_local_endpoint_rejects_non_loopback_address(self):
        with self.assertRaises(ValueError):
            llm.use_local_endpoint('http://192.168.1.10:8080/v1')

    def test_local_endpoint_detection_requires_loopback(self):
        with patch.object(llm, '_config', {'base_url': 'http://example.test/v1'}):
            self.assertFalse(llm.is_local_endpoint())
        with patch.object(llm, '_config', {'base_url': 'http://127.0.0.1:8080/v1'}):
            self.assertTrue(llm.is_local_endpoint())

    def test_task_model_overrides_mode_then_falls_back_to_default(self):
        with patch.object(llm, '_config', {
            'model': 'default-model',
            'model_mode': 'fast',
            'model_modes': {'fast': 'fast-model'},
            'task_models': {'translation': 'translation-model'},
        }):
            self.assertEqual(llm.model_for_task('translation'), 'translation-model')
            self.assertEqual(llm.model_for_task('chat'), 'fast-model')
        with patch.object(llm, '_config', {'model': 'default-model'}):
            self.assertEqual(llm.model_for_task('document'), 'default-model')

    def test_public_settings_never_returns_api_key(self):
        with patch.object(llm, '_config', {'api_key': 'secret', 'model': 'model-a'}):
            self.assertNotIn('api_key', llm.public_settings())

    def test_toolbox_settings_use_atomic_config_merge_and_validate_url(self):
        with tempfile.TemporaryDirectory() as directory:
            config_path = Path(directory) / 'config.json'
            with patch.object(settings, 'CONFIG_FILE', config_path), \
                 patch.object(settings, 'config', {}), \
                 patch.object(llm, '_config', {}):
                saved = llm.save_toolbox_settings({
                    'provider': 'company', 'base_url': 'http://intranet.test/v1',
                    'model': 'default-model', 'request_timeout': 90,
                    'model_modes': {'fast': 'small-model'},
                    'task_models': {'translation': 'translate-model'},
                })
                data = json.loads(config_path.read_text(encoding='utf-8'))
            self.assertEqual(saved['request_timeout'], 90)
            self.assertEqual(data['llm']['task_models']['translation'], 'translate-model')

    def test_blank_timeout_does_not_discard_other_settings(self):
        with tempfile.TemporaryDirectory() as directory:
            config_path = Path(directory) / 'config.json'
            with patch.object(settings, 'CONFIG_FILE', config_path), \
                 patch.object(settings, 'config', {}), \
                 patch.object(llm, '_config', {}):
                llm.save_toolbox_settings({'model': 'model-a', 'request_timeout': ''})
            data = json.loads(config_path.read_text(encoding='utf-8'))
        self.assertEqual(data['llm']['model'], 'model-a')

    def test_sidecar_model_overrides_remote_task_routes(self):
        with patch.object(llm, '_config', {
            'base_url': 'http://company.test/v1', 'model': 'company-model',
            'task_models': {'document': 'company-document-model'},
        }), patch.object(llm, '_local_sidecar_active', False):
            llm.use_local_endpoint('http://127.0.0.1:8080/v1', 'local-alias')
            self.assertEqual(llm.model_for_task('document'), 'local-alias')


class CodexLimitsTests(unittest.TestCase):
    def test_millisecond_reset_timestamp_is_normalized(self):
        self.assertEqual(
            codex_limits._unix_to_iso(1_700_000_000_000),
            codex_limits._unix_to_iso(1_700_000_000),
        )

    def test_rpc_error_keeps_safe_server_diagnostic(self):
        message = codex_limits._rpc_error_message({
            'code': -32601, 'message': 'Method not found',
        })
        self.assertIn('-32601', message)
        self.assertIn('Method not found', message)


class JsonToolTests(unittest.TestCase):
    def test_format_preserves_object_key_order_without_llm(self):
        result = json_tools.process('{"z":1,"a":[true,null]}', 'format')
        self.assertTrue(result['valid'])
        self.assertEqual(result['text'], '{\n  "z": 1,\n  "a": [\n    true,\n    null\n  ]\n}')

    def test_invalid_json_reports_line_and_column(self):
        result = json_tools.process('{\n  "a":\n}', 'validate')
        self.assertFalse(result['valid'])
        self.assertEqual(result['line'], 3)
        self.assertEqual(result['column'], 1)
        self.assertIn('第 3 行、第', result['error'])

    def test_search_reports_jsonpath_for_key_and_value(self):
        result = json_tools.process('{"data":[{"id":"PO-1"}]}', 'search', 'po-1')
        self.assertEqual(result['matches'][0]['path'], '$["data"][0]["id"]')

    def test_excessive_json_depth_is_rejected_before_rendering(self):
        text = '[' * 101 + '0' + ']' * 101
        result = json_tools.process(text, 'format')
        self.assertFalse(result['valid'])
        self.assertIn('巢狀層級', result['error'])


class TranslationTests(unittest.TestCase):
    def test_technical_translation_prompt_protects_identifiers_and_uses_glossary(self):
        with patch.object(llm, '_config', {'translation_glossary': {'Receipt': '收料'}}):
            messages = translation_service.build_messages('SELECT ORG_ID FROM po', {
                'target': 'zh-TW', 'mode': 'technical', 'preserve_code': True,
                'preserve_tables': True, 'use_glossary': True,
            })
        system = messages[0]['content']
        self.assertIn('SQL 欄位名稱', system)
        self.assertIn('JSON Key', system)
        self.assertIn('Receipt → 收料', system)

    def test_simplified_chinese_prompt_and_same_language_rejection(self):
        messages = translation_service.build_messages('測試', {
            'source': 'zh-TW', 'target': 'zh-CN',
        })
        system = messages[0]['content']
        self.assertIn('來源語言：繁體中文', system)
        self.assertIn('目標語言：簡體中文', system)
        self.assertIn('簡體中文字形', system)
        with patch.object(llm, 'chat') as chat:
            result = translation_service.translate('測試', {
                'source': 'zh-CN', 'target': 'zh-CN',
            })
        self.assertEqual(result['error'], '來源語言與目標語言不可相同。')
        chat.assert_not_called()

    def test_translation_uses_task_routed_model(self):
        with patch.object(llm, 'model_for_task', return_value='translate-model'), \
             patch.object(llm, 'request_timeout', return_value=90), \
             patch.object(llm, 'chat', return_value={'content': '譯文', 'model': 'translate-model'}) as chat:
            result = translation_service.translate('hello', {'target': 'zh-TW'})
        self.assertEqual(result['content'], '譯文')
        self.assertEqual(chat.call_args.kwargs['model'], 'translate-model')
        self.assertEqual(chat.call_args.kwargs['timeout'], 90)

    def test_translation_restores_protected_sql_and_rejects_lost_placeholder(self):
        source = "請執行 SELECT ORG_ID FROM po WHERE status = 'OPEN'。"
        protected, placeholders = translation_service.protect_content(source, {'preserve_code': True})
        with patch.object(llm, 'chat', return_value={'content': protected}), \
             patch.object(llm, 'model_for_task', return_value='translate-model'):
            result = translation_service.translate(source, {'preserve_code': True})
        self.assertIn('SELECT', result['content'])
        with self.assertRaises(ValueError):
            translation_service.restore_content('placeholder lost', placeholders)

    def test_translation_reaches_openai_compatible_endpoint_with_task_model(self):
        received = []

        class Handler(BaseHTTPRequestHandler):
            def do_POST(self):
                length = int(self.headers['Content-Length'])
                received.append(json.loads(self.rfile.read(length)))
                body = json.dumps({'model': received[-1]['model'], 'choices': [
                    {'message': {'content': 'translated'}}]}).encode('utf-8')
                self.send_response(200)
                self.send_header('Content-Type', 'application/json')
                self.send_header('Content-Length', str(len(body)))
                self.end_headers()
                self.wfile.write(body)

            def log_message(self, _format, *_args):
                return

        server = ThreadingHTTPServer(('127.0.0.1', 0), Handler)
        thread = threading.Thread(target=server.serve_forever, daemon=True)
        thread.start()
        try:
            with patch.object(llm, '_config', {
                'base_url': f'http://127.0.0.1:{server.server_address[1]}/v1',
                'model': 'chat-model', 'task_models': {'translation': 'translate-model'},
            }):
                result = translation_service.translate('hello', {'preserve_code': False})
        finally:
            server.shutdown()
            thread.join(timeout=2)
            server.server_close()
        self.assertEqual(result['content'], 'translated')
        self.assertEqual(received[0]['model'], 'translate-model')


class SettingsTests(unittest.TestCase):
    def test_merge_config_preserves_independent_sections(self):
        with tempfile.TemporaryDirectory() as directory:
            config_path = Path(directory) / 'config.json'
            with patch.object(settings, 'CONFIG_FILE', config_path), \
                 patch.object(settings, 'config', {}):
                settings.merge_config({'llm': {'model': 'model-a'}})
                settings.merge_config({'skin': 'bluecat', 'llm': {'debug_log': True}})
                data = json.loads(config_path.read_text(encoding='utf-8'))
            self.assertEqual(data['skin'], 'bluecat')
            self.assertEqual(data['llm'], {'model': 'model-a', 'debug_log': True})



class LocalLlmRuntimeTests(unittest.TestCase):
    def test_disabled_runtime_does_not_start_a_process(self):
        with tempfile.TemporaryDirectory() as directory:
            config = Path(directory) / 'config.json'
            config.write_text(json.dumps({'local_llm': {'enabled': False}}), encoding='utf-8')
            self.assertEqual(local_llm.init(config)['status'], '本機模型未啟用。')

    def test_runtime_reports_server_start_failure(self):
        with tempfile.TemporaryDirectory() as directory:
            root = Path(directory)
            config = root / 'config.json'
            (root / 'llama-server.exe').touch()
            (root / 'model.gguf').touch()
            config.write_text(json.dumps({'local_llm': {'enabled': True}}), encoding='utf-8')
            with patch.object(local_llm.subprocess, 'Popen', side_effect=OSError('blocked')):
                result = local_llm.init(config)
            self.assertIn('無法啟動本機模型服務', result['status'])
            self.assertIsNone(result['endpoint'])


class DocumentServiceTests(unittest.TestCase):
    def test_markdown_query_returns_deterministic_line_citation(self):
        with tempfile.TemporaryDirectory() as directory:
            root = Path(directory)
            source = root / '採購.md'
            source.write_text('# 付款條款\n付款期限為收貨後 30 日。\n', encoding='utf-8')
            with patch.object(documents, 'DOCUMENTS_DIR', root / 'documents'):
                result = documents.ingest(str(source))
                self.assertNotIn('error', result)
                answer = documents.query(result['document']['id'], '付款期限')
                self.assertEqual(answer['answer'], '找到下列與問題相關的文件內容：')
                self.assertEqual(answer['sources'][0]['source']['document_name'], '採購.md')
                self.assertEqual(answer['sources'][0]['source']['line_start'], 2)
                self.assertEqual(answer['sources'][0]['source']['heading'], '付款條款')

    def test_missing_evidence_does_not_invent_an_answer(self):
        with tempfile.TemporaryDirectory() as directory:
            root = Path(directory)
            source = root / '流程.txt'
            source.write_text('申請人填寫表單。\n', encoding='utf-8')
            with patch.object(documents, 'DOCUMENTS_DIR', root / 'documents'):
                result = documents.ingest(str(source))
                answer = documents.query(result['document']['id'], '付款期限')
                self.assertEqual(answer['answer'], '此文件沒有描述此問題，無法依文件確認。')
                self.assertEqual(answer['sources'], [])

    def test_chinese_common_words_do_not_make_unrelated_policy_evidence(self):
        with tempfile.TemporaryDirectory() as directory:
            root = Path(directory)
            source = root / 'leave.md'
            source.write_text('# 年假\n年假可以請 30 天。\n', encoding='utf-8')
            with patch.object(documents, 'DOCUMENTS_DIR', root / 'documents'):
                result = documents.ingest(str(source))
                answer = documents.query(result['document']['id'], '育嬰假可以請幾天？')
            self.assertEqual(answer['sources'], [])

    def test_document_context_samples_and_discloses_partial_coverage(self):
        with tempfile.TemporaryDirectory() as directory:
            root = Path(directory)
            source = root / 'long.md'
            source.write_text('\n\n'.join(f'第 {index} 段內容' for index in range(20)), encoding='utf-8')
            with patch.object(documents, 'DOCUMENTS_DIR', root / 'documents'):
                result = documents.ingest(str(source))
                evidence = documents.context(result['document']['id'], limit=3)
            self.assertFalse(evidence['coverage']['complete'])
            self.assertEqual(evidence['coverage']['included_chunks'], 3)
            self.assertEqual(evidence['sources'][-1]['excerpt'], '第 19 段內容')

    def test_scanned_pdf_is_rejected_with_an_ocr_message(self):
        from pypdf import PdfWriter
        with tempfile.TemporaryDirectory() as directory:
            source = Path(directory) / 'scan.pdf'
            writer = PdfWriter()
            writer.add_blank_page(width=100, height=100)
            with source.open('wb') as output:
                writer.write(output)
            result = documents.ingest(str(source))
            self.assertEqual(result['error'], '此掃描型 PDF 需先經 OCR 才能閱讀。')

    def test_xlsx_query_returns_sheet_and_cell_range(self):
        from openpyxl import Workbook
        with tempfile.TemporaryDirectory() as directory:
            root = Path(directory)
            source = root / '付款.xlsx'
            workbook = Workbook()
            sheet = workbook.active
            sheet.title = '付款條款'
            sheet.append(['條件', '內容'])
            sheet.append(['付款期限', '收貨後 30 日'])
            workbook.save(source)
            with patch.object(documents, 'DOCUMENTS_DIR', root / 'documents'):
                result = documents.ingest(str(source))
                answer = documents.query(result['document']['id'], '付款期限')
                self.assertEqual(answer['sources'][0]['source']['sheet'], '付款條款')
                self.assertEqual(answer['sources'][0]['source']['cell_range'], 'A2:B2')

    def test_xlsx_evidence_keeps_header_with_value_row(self):
        from openpyxl import Workbook
        with tempfile.TemporaryDirectory() as directory:
            root = Path(directory)
            source = root / 'terms.xlsx'
            workbook = Workbook()
            sheet = workbook.active
            sheet.append(['條件', '內容'])
            sheet.append(['付款期限', '收貨後 30 日'])
            workbook.save(source)
            with patch.object(documents, 'DOCUMENTS_DIR', root / 'documents'):
                result = documents.ingest(str(source))
                answer = documents.query(result['document']['id'], '付款期限')
            self.assertIn('內容: 收貨後 30 日', answer['sources'][0]['excerpt'])

    def test_office_and_pdf_extractors_preserve_native_locations(self):
        from docx import Document
        from pptx import Presentation
        from reportlab.pdfgen.canvas import Canvas
        with tempfile.TemporaryDirectory() as directory:
            root = Path(directory)
            word_path = root / 'policy.docx'
            word = Document()
            word.add_heading('Payment terms', level=1)
            word.add_paragraph('Payment is due in 30 days.')
            word.save(word_path)

            ppt_path = root / 'policy.pptx'
            deck = Presentation()
            slide = deck.slides.add_slide(deck.slide_layouts[1])
            slide.shapes.title.text = 'Approval'
            slide.placeholders[1].text = 'Manager approval is required.'
            deck.save(ppt_path)

            pdf_path = root / 'policy.pdf'
            pdf = Canvas(str(pdf_path))
            pdf.drawString(72, 720, 'Payment is due in 30 days.')
            pdf.save()

            with patch.object(documents, 'DOCUMENTS_DIR', root / 'documents'):
                word_result = documents.ingest(str(word_path))
                ppt_result = documents.ingest(str(ppt_path))
                pdf_result = documents.ingest(str(pdf_path))
                self.assertEqual(documents.context(word_result['document']['id'])['sources'][0]['source']['kind'], 'word_paragraph')
                self.assertEqual(documents.context(ppt_result['document']['id'])['sources'][0]['source']['locator'], '投影片 1 · Approval')
                self.assertEqual(documents.context(pdf_result['document']['id'])['sources'][0]['source']['locator'], '第 1 頁')

    def test_document_question_sends_only_retrieved_evidence_to_llm(self):
        evidence = {
            'answer': '找到下列與問題相關的文件內容：',
            'sources': [{'excerpt': '付款期限為收貨後 30 日。', 'source': {
                'document_name': '採購.md', 'locator': '第 2 行 · 付款條款'}}],
        }
        with patch.object(documents, 'query', return_value=evidence), \
             patch.object(llm, 'chat', return_value={'content': '付款期限為收貨後 30 日。'} ) as chat:
            result = JsApi().query_document('550e8400-e29b-41d4-a716-446655440000', '付款期限？')
        self.assertEqual(result['answer'], '付款期限為收貨後 30 日。')
        prompt = chat.call_args.args[0][1]['content']
        self.assertIn('付款期限為收貨後 30 日。', prompt)
        self.assertIn('採購.md · 第 2 行 · 付款條款', prompt)

    def test_document_summary_action_carries_source_metadata(self):
        evidence = {'sources': [{'excerpt': '採購前需取得主管核准。', 'source': {
            'document_name': '流程.md', 'locator': '第 4 行'}}]}
        with patch.object(documents, 'context', return_value=evidence), \
             patch.object(llm, 'chat', return_value={'content': '需先取得主管核准。'}) as chat:
            result = JsApi().document_action('550e8400-e29b-41d4-a716-446655440000', 'summary')
        self.assertEqual(result['answer'], '需先取得主管核准。')
        self.assertEqual(result['sources'], evidence['sources'])
        self.assertIn('流程.md · 第 4 行', chat.call_args.args[0][1]['content'])

    def test_document_comparison_carries_both_documents_sources(self):
        first = {'sources': [{'excerpt': '付款 30 日。', 'source': {'document_name': 'A.md', 'locator': '第 1 行'}}]}
        second = {'sources': [{'excerpt': '付款 45 日。', 'source': {'document_name': 'B.md', 'locator': '第 1 行'}}]}
        with patch.object(documents, 'context', side_effect=[first, second]), \
             patch.object(llm, 'chat', return_value={'content': '兩份文件付款期限不同。'}) as chat:
            result = JsApi().compare_documents('550e8400-e29b-41d4-a716-446655440000', '11111111-1111-4111-8111-111111111111')
        self.assertEqual(len(result['sources']), 2)
        prompt = chat.call_args.args[0][1]['content']
        self.assertIn('A.md · 第 1 行', prompt)
        self.assertIn('B.md · 第 1 行', prompt)


class WorkflowTests(unittest.TestCase):
    DOCUMENT_ID = '550e8400-e29b-41d4-a716-446655440000'
    EVIDENCE = {
        'sources': [{
            'excerpt': '付款期限為收貨後 30 日。',
            'source': {'document_name': '採購流程.pdf', 'locator': '第 8 頁'},
        }],
        'coverage': {'included_chunks': 1, 'total_chunks': 2, 'complete': False},
    }

    def _directories(self, root):
        return (
            patch.object(workflows, 'RUNS_DIR', root / 'runs'),
            patch.object(workflows, 'ARTIFACTS_DIR', root / 'artifacts'),
        )

    def test_document_meeting_pack_completes_with_source_bearing_artifact(self):
        with tempfile.TemporaryDirectory() as directory:
            root = Path(directory)
            runs_patch, artifacts_patch = self._directories(root)
            with runs_patch, artifacts_patch, \
                 patch.object(documents, 'context', return_value=self.EVIDENCE), \
                 patch.object(llm, 'chat', side_effect=[
                     {'content': '這是抽樣摘要。', 'model': 'company-doc'},
                     {'content': '決策：確認付款期限。', 'model': 'company-doc'},
                 ]):
                created = workflows.create_document_meeting_pack(self.DOCUMENT_ID)
                result = workflows.execute(created['run_id'])

            self.assertEqual(result['status'], 'completed')
            self.assertTrue(all(step['status'] == 'completed' for step in result['steps']))
            artifact = Path(result['artifacts'][0]['path'])
            self.assertTrue(artifact.is_file())
            text = artifact.read_text(encoding='utf-8')
            self.assertIn('這是抽樣摘要。', text)
            self.assertIn('決策：確認付款期限。', text)
            self.assertIn('採購流程.pdf · 第 8 頁', text)
            self.assertIn('抽樣涵蓋 1/2', text)
            self.assertEqual(result['model'], 'company-doc')
            self.assertEqual(result['coverage'], self.EVIDENCE['coverage'])
            self.assertEqual(result['sources'], [{
                'document_name': '採購流程.pdf',
                'locator': '第 8 頁',
            }])

    def test_later_llm_failure_preserves_completed_summary_as_partial_artifact(self):
        with tempfile.TemporaryDirectory() as directory:
            root = Path(directory)
            runs_patch, artifacts_patch = self._directories(root)
            with runs_patch, artifacts_patch, \
                 patch.object(documents, 'context', return_value=self.EVIDENCE), \
                 patch.object(llm, 'chat', side_effect=[
                     {'content': '已完成摘要。', 'model': 'company-doc'},
                     {'error': '會議重點模型逾時'},
                 ]):
                created = workflows.create_document_meeting_pack(self.DOCUMENT_ID)
                result = workflows.execute(created['run_id'])

            self.assertEqual(result['status'], 'failed')
            self.assertEqual(result['steps'][1]['status'], 'completed')
            self.assertEqual(result['steps'][2]['status'], 'failed')
            self.assertEqual(result['artifacts'][0]['status'], 'partial')
            self.assertIn('已完成摘要。',
                          Path(result['artifacts'][0]['path']).read_text(encoding='utf-8'))

    def test_export_failure_never_marks_workflow_completed(self):
        def fail_export(_run, _data):
            raise OSError('磁碟已滿')

        with tempfile.TemporaryDirectory() as directory:
            root = Path(directory)
            runs_patch, artifacts_patch = self._directories(root)
            with runs_patch, artifacts_patch, \
                 patch.object(documents, 'context', return_value=self.EVIDENCE), \
                 patch.object(llm, 'chat', side_effect=[
                     {'content': '摘要。', 'model': 'company-doc'},
                     {'content': '會議重點。', 'model': 'company-doc'},
                 ]), patch.dict(workflows._HANDLERS, {'export_markdown': fail_export}):
                created = workflows.create_document_meeting_pack(self.DOCUMENT_ID)
                result = workflows.execute(created['run_id'])

            self.assertEqual(result['status'], 'failed')
            self.assertEqual(result['steps'][-1]['status'], 'failed')
            self.assertIn('磁碟已滿', result['steps'][-1]['error'])

    def test_cancelled_pending_run_does_not_call_document_or_llm(self):
        with tempfile.TemporaryDirectory() as directory:
            root = Path(directory)
            runs_patch, artifacts_patch = self._directories(root)
            with runs_patch, artifacts_patch, \
                 patch.object(documents, 'context') as context, \
                 patch.object(llm, 'chat') as chat:
                created = workflows.create_document_meeting_pack(self.DOCUMENT_ID)
                workflows.cancel_run(created['run_id'])
                result = workflows.execute(created['run_id'])

            self.assertEqual(result['status'], 'cancelled')
            context.assert_not_called()
            chat.assert_not_called()

    def test_document_meeting_pack_rejects_non_pdf_docx_input(self):
        evidence = {
            **self.EVIDENCE,
            'sources': [{
                'excerpt': '文字內容。',
                'source': {'document_name': 'notes.txt', 'locator': '第 1 行'},
            }],
        }
        with tempfile.TemporaryDirectory() as directory:
            root = Path(directory)
            runs_patch, artifacts_patch = self._directories(root)
            with runs_patch, artifacts_patch, \
                 patch.object(documents, 'context', return_value=evidence), \
                 patch.object(llm, 'chat') as chat:
                created = workflows.create_document_meeting_pack(self.DOCUMENT_ID)
                result = workflows.execute(created['run_id'])

            self.assertEqual(result['status'], 'failed')
            self.assertEqual(result['steps'][0]['status'], 'failed')
            self.assertIn('只支援 PDF 與 DOCX', result['steps'][0]['error'])
            chat.assert_not_called()

    def test_failed_workflow_retry_creates_a_new_bounded_run(self):
        with tempfile.TemporaryDirectory() as directory:
            root = Path(directory)
            runs_patch, artifacts_patch = self._directories(root)
            with runs_patch, artifacts_patch:
                previous = workflows.create_document_meeting_pack(self.DOCUMENT_ID, translate=True)
                previous['status'] = 'failed'
                workflows._write_run(previous)
                retried = workflows.retry_run(previous['run_id'])

                self.assertNotEqual(retried['run_id'], previous['run_id'])
                self.assertEqual(retried['retry_of'], previous['run_id'])
                self.assertEqual(retried['retry_count'], 1)
                self.assertIn('translate', [step['id'] for step in retried['steps']])
                self.assertEqual(workflows.latest_run()['run_id'], retried['run_id'])
                duplicate = workflows.retry_run(previous['run_id'])
                self.assertIn('已經建立過', duplicate['error'])

                retried['status'] = 'failed'
                retried['retry_count'] = workflows.MAX_RETRIES
                workflows._write_run(retried)
                blocked = workflows.retry_run(retried['run_id'])
                self.assertIn('重新執行上限', blocked['error'])

    def test_completed_workflow_cannot_be_retried(self):
        with tempfile.TemporaryDirectory() as directory:
            root = Path(directory)
            runs_patch, artifacts_patch = self._directories(root)
            with runs_patch, artifacts_patch:
                previous = workflows.create_document_meeting_pack(self.DOCUMENT_ID)
                previous['status'] = 'completed'
                workflows._write_run(previous)
                result = workflows.retry_run(previous['run_id'])
            self.assertIn('只有失敗或已取消', result['error'])

    def test_retry_repairs_pointer_to_missing_successor(self):
        with tempfile.TemporaryDirectory() as directory:
            root = Path(directory)
            runs_patch, artifacts_patch = self._directories(root)
            with runs_patch, artifacts_patch:
                previous = workflows.create_document_meeting_pack(self.DOCUMENT_ID)
                previous['status'] = 'failed'
                previous['retry_to'] = '11111111-1111-4111-8111-111111111111'
                workflows._write_run(previous)
                result = workflows.retry_run(previous['run_id'])
            self.assertEqual(result['status'], 'pending')
            self.assertEqual(result['retry_of'], previous['run_id'])

    def test_latest_run_skips_newest_corrupt_json(self):
        with tempfile.TemporaryDirectory() as directory:
            root = Path(directory)
            runs_patch, artifacts_patch = self._directories(root)
            with runs_patch, artifacts_patch:
                valid = workflows.create_document_meeting_pack(self.DOCUMENT_ID)
                valid['status'] = 'completed'
                workflows._write_run(valid)
                corrupt = root / 'runs' / 'ffffffff-ffff-4fff-8fff-ffffffffffff.json'
                corrupt.write_text('{broken', encoding='utf-8')
                result = workflows.latest_run()
            self.assertEqual(result['run_id'], valid['run_id'])

    def test_retention_and_manual_cleanup_bound_local_history(self):
        with tempfile.TemporaryDirectory() as directory:
            root = Path(directory)
            runs_patch, artifacts_patch = self._directories(root)
            with runs_patch, artifacts_patch, patch.object(workflows, 'MAX_COMPLETED_RUNS', 2):
                for index in range(4):
                    run = workflows.create_document_meeting_pack(self.DOCUMENT_ID)
                    run['status'] = 'completed'
                    workflows._write_run(run)
                    artifact = root / 'artifacts' / f'{run["run_id"]}_meeting-pack.md'
                    artifact.parent.mkdir(parents=True, exist_ok=True)
                    artifact.write_text(f'run {index}', encoding='utf-8')
                pruned = workflows.prune_history()
                self.assertGreaterEqual(pruned['removed_runs'], 1)
                self.assertLessEqual(len(list((root / 'runs').glob('*.json'))), 2)
                self.assertLessEqual(len(list((root / 'artifacts').glob('*.md'))), 2)

                active = workflows.create_document_meeting_pack(self.DOCUMENT_ID)
                active_artifact = root / 'artifacts' / f'{active["run_id"]}_meeting-pack.md'
                active_artifact.write_text('partial', encoding='utf-8')
                cleared = workflows.clear_history()
                self.assertEqual(cleared['active_runs_preserved'], 1)
                self.assertTrue(workflows._run_path(active['run_id']).exists())
                self.assertTrue(active_artifact.exists())

    def test_stale_running_run_is_recovered_as_retryable_failure(self):
        with tempfile.TemporaryDirectory() as directory:
            root = Path(directory)
            runs_patch, artifacts_patch = self._directories(root)
            with runs_patch, artifacts_patch:
                stale = workflows.create_document_meeting_pack(self.DOCUMENT_ID)
                stale['status'] = 'running'
                stale['owner_session'] = 'previous-process'
                stale['current_step'] = stale['steps'][0]['id']
                stale['steps'][0]['status'] = 'running'
                workflows._write_run(stale)

                recovered = workflows.latest_run()
                self.assertEqual(recovered['status'], 'failed')
                self.assertIn('上次程式中斷', recovered['recovery_error'])
                self.assertEqual(recovered['steps'][0]['status'], 'failed')
                retried = workflows.retry_run(recovered['run_id'])
                self.assertEqual(retried['status'], 'pending')
                another = workflows.create_document_meeting_pack(self.DOCUMENT_ID)
                another['status'] = 'running'
                another['owner_session'] = 'previous-process'
                workflows._write_run(another)
                cleared = workflows.clear_history()
                self.assertGreaterEqual(cleared['removed_runs'], 1)
                self.assertFalse(workflows._run_path(another['run_id']).exists())
                self.assertTrue(workflows._run_path(retried['run_id']).exists())

    def test_retained_corrupt_run_keeps_its_markdown_artifact(self):
        with tempfile.TemporaryDirectory() as directory:
            root = Path(directory)
            runs_patch, artifacts_patch = self._directories(root)
            with runs_patch, artifacts_patch:
                run = workflows.create_document_meeting_pack(self.DOCUMENT_ID)
                artifact = root / 'artifacts' / f'{run["run_id"]}_meeting-pack.md'
                artifact.parent.mkdir(parents=True, exist_ok=True)
                artifact.write_text('deliverable', encoding='utf-8')
                workflows._run_path(run['run_id']).write_text('{broken', encoding='utf-8')

                workflows.prune_history()
                self.assertTrue(workflows._run_path(run['run_id']).exists())
                self.assertTrue(artifact.exists())


class WorkerTests(unittest.TestCase):
    def test_xlsx_worker_uses_openpyxl_without_pandas(self):
        from openpyxl import Workbook
        with tempfile.TemporaryDirectory() as directory:
            source = Path(directory) / 'sample.xlsx'
            workbook = Workbook()
            workbook.active.append(['name', 'value'])
            workbook.active.append(['payment', 30])
            workbook.save(source)
            workbook.close()
            output = io.StringIO()
            with redirect_stdout(output):
                worker.main(str(source))
            gc.collect()
            self.assertIn('payment,30', output.getvalue())

    def test_xlsx_worker_reads_all_sheets_and_reports_row_truncation(self):
        from openpyxl import Workbook
        with tempfile.TemporaryDirectory() as directory:
            source = Path(directory) / 'sheets.xlsx'
            workbook = Workbook()
            workbook.active.title = 'Summary'
            workbook.active.append(['summary'])
            detail = workbook.create_sheet('Detail')
            detail.append(['row-1'])
            detail.append(['row-2'])
            detail.append(['row-3'])
            workbook.save(source)
            workbook.close()
            output = io.StringIO()
            with patch.object(worker, 'MAX_ROWS_PER_SHEET', 2), redirect_stdout(output):
                worker.main(str(source))
            gc.collect()
            text = output.getvalue()
        self.assertIn('# 工作表：Summary', text)
        self.assertIn('# 工作表：Detail', text)
        self.assertIn('僅提供前 2 列', text)

    def test_ppt_worker_accepts_h2_slide_headings(self):
        from pptx import Presentation
        with tempfile.TemporaryDirectory() as directory:
            target = Path(directory) / 'slides.pptx'
            source = '## Slide 1\n- First\n\n## Slide 2\n- Second\n'
            with patch.object(worker.sys, 'stdin', io.StringIO(source)), redirect_stdout(io.StringIO()):
                worker.generate_ppt(str(target))
            deck = Presentation(target)
        self.assertEqual(len(deck.slides), 2)
        self.assertEqual(deck.slides[1].shapes.title.text, '2')


if __name__ == '__main__':
    unittest.main()
