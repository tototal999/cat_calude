"""End-to-end verifier for the packaged document indexers.

Run with the same Python used to build the EXE. It creates four temporary
documents, invokes the packaged executable for each, and checks a real source
locator through the ``--document-check`` entry point.
"""
from __future__ import annotations

import os
import subprocess
import sys
import tempfile
from pathlib import Path


def _make_samples(root: Path) -> list[tuple[Path, str, str]]:
    from docx import Document
    from openpyxl import Workbook
    from pptx import Presentation
    from reportlab.pdfgen.canvas import Canvas

    pdf_path = root / 'payment.pdf'
    pdf = Canvas(str(pdf_path))
    pdf.drawString(72, 720, 'Payment is due in 30 days.')
    pdf.save()

    docx_path = root / 'payment.docx'
    word = Document()
    word.add_heading('Payment terms', level=1)
    word.add_paragraph('Payment is due in 30 days.')
    word.save(docx_path)

    pptx_path = root / 'approval.pptx'
    deck = Presentation()
    slide = deck.slides.add_slide(deck.slide_layouts[1])
    slide.shapes.title.text = 'Approval'
    slide.placeholders[1].text = 'Manager approval is required.'
    deck.save(pptx_path)

    xlsx_path = root / 'payment.xlsx'
    workbook = Workbook()
    sheet = workbook.active
    sheet.title = 'Payment'
    sheet.append(['Term', 'Value'])
    sheet.append(['Payment due', '30 days'])
    workbook.save(xlsx_path)

    return [
        (pdf_path, 'Payment', '第 1 頁'),
        (docx_path, 'Payment', '第 2 段 · Payment terms'),
        (pptx_path, 'approval', '投影片 1 · Approval'),
        (xlsx_path, 'Payment', '工作表 Payment · A2:B2'),
    ]


def main() -> int:
    executable = Path(sys.argv[1]) if len(sys.argv) > 1 else Path('dist/ClaudeCat/ClaudeCat.exe')
    if not executable.is_file():
        print(f'QA_RESULT|STATUS:FAIL|EXPECTED:packaged executable|ACTUAL:missing {executable}')
        return 1
    with tempfile.TemporaryDirectory() as directory:
        root = Path(directory)
        environment = os.environ.copy()
        environment['LOCALAPPDATA'] = str(root / 'appdata')
        for path, query, locator in _make_samples(root):
            completed = subprocess.run(
                [str(executable), '--document-check', str(path), query, locator],
                capture_output=True,
                text=True,
                timeout=60,
                env=environment,
            )
            if completed.returncode != 0:
                print(f'QA_RESULT|STATUS:FAIL|EXPECTED:{path.name} -> {locator}|ACTUAL:exit {completed.returncode}')
                return 1
    print('QA_RESULT|STATUS:PASS|EXPECTED:PDF/DOCX/PPTX/XLSX index and citation|ACTUAL:all formats passed')
    return 0


if __name__ == '__main__':
    raise SystemExit(main())
