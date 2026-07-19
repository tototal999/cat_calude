import sys
from pathlib import Path
import json
import csv
import re
from io import StringIO

MAX_ROWS_PER_SHEET = 10_000

for _stream in (sys.stdout, sys.stderr):
    if _stream is not None and hasattr(_stream, 'reconfigure'):
        _stream.reconfigure(encoding='utf-8', errors='replace')

def main(filepath_str, max_chars=50000):
    try:
        p = Path(filepath_str)
        if not p.exists():
            print(f"檔案不存在: {p}", file=sys.stderr)
            sys.exit(1)
            
        if p.suffix.lower() == '.xlsx':
            from openpyxl import load_workbook
            workbook = load_workbook(p, read_only=False, data_only=True)
            try:
                output = StringIO()
                writer = csv.writer(output)
                truncated = []
                for sheet in workbook.worksheets:
                    writer.writerow([f'# 工作表：{sheet.title}'])
                    rows = sheet.iter_rows(values_only=True)
                    try:
                        for row_no, row in enumerate(rows, start=1):
                            if row_no > MAX_ROWS_PER_SHEET:
                                truncated.append(sheet.title)
                                break
                            writer.writerow(['' if value is None else value for value in row])
                    finally:
                        rows.close()
                if truncated:
                    writer.writerow([f'# 警告：{", ".join(truncated)} 僅提供前 {MAX_ROWS_PER_SHEET:,} 列。'])
                content = output.getvalue()
            finally:
                workbook.close()
        elif p.suffix.lower() == '.xls':
            import xlrd
            workbook = xlrd.open_workbook(p, on_demand=True)
            try:
                output = StringIO()
                writer = csv.writer(output)
                truncated = []
                for sheet_no in range(workbook.nsheets):
                    sheet = workbook.sheet_by_index(sheet_no)
                    writer.writerow([f'# 工作表：{sheet.name}'])
                    for row_no in range(min(sheet.nrows, MAX_ROWS_PER_SHEET)):
                        writer.writerow(sheet.row_values(row_no))
                    if sheet.nrows > MAX_ROWS_PER_SHEET:
                        truncated.append(sheet.name)
                if truncated:
                    writer.writerow([f'# 警告：{", ".join(truncated)} 僅提供前 {MAX_ROWS_PER_SHEET:,} 列。'])
                content = output.getvalue()
            finally:
                workbook.release_resources()
        else:
            with p.open('r', encoding='utf-8') as source:
                content = source.read(max_chars + 1)
        if len(content) > max_chars:
            raise ValueError(f'File content exceeds the {max_chars}-character limit')
        print(content, end='')
    except Exception as e:
        print(f"Worker 解析失敗: {e}", file=sys.stderr)
        sys.exit(1)

def generate_ppt(output_path_str, template_path_str=None):
    try:
        from pptx import Presentation
        from pptx.util import Pt
        
        import os
        if sys.stdin is None:
            # PyInstaller windowed mode fallback
            content = b""
            try:
                while True:
                    chunk = os.read(0, 8192)
                    if not chunk: break
                    content += chunk
            except Exception:
                pass
            content = content.decode('utf-8')
        else:
            content = sys.stdin.read()
        
        if template_path_str:
            prs = Presentation(template_path_str)
        else:
            prs = Presentation()
        
        # Each top-level Markdown heading is a slide.  Some models emit
        # ``## Slide 2`` instead; accept that form when no H1 slide exists.
        headings = list(re.finditer(r'(?m)^#\s+(.+?)\s*$', content))
        if not headings:
            headings = list(re.finditer(r'(?mi)^##\s+(slide\b.+?)\s*$', content))
        if not headings:
            raise ValueError('簡報內容需以「# 標題」或「## Slide 標題」分頁。')
        slides_data = []
        for index, heading in enumerate(headings):
            end = headings[index + 1].start() if index + 1 < len(headings) else len(content)
            title = re.sub(r'(?i)^slide\s*:?\s*', '', heading.group(1)).strip()
            slides_data.append((title, content[heading.end():end]))
        
        for title, slide_text in slides_data:
            lines = slide_text.strip().split('\n')
            
            # Find bullet points
            bullets = []
            notes = []
            in_notes = False
            for line in lines[1:]:
                line = line.strip()
                if not line:
                    continue
                if line.lower().startswith('## notes') or line.lower().startswith('## 備註') or line.lower().startswith('## 講者備註'):
                    in_notes = True
                    continue
                if line.startswith('##'):
                    in_notes = False
                    continue
                    
                if in_notes:
                    notes.append(line)
                elif line.startswith('- ') or line.startswith('* '):
                    bullets.append(line[2:].strip())
                elif line:
                    bullets.append(line)
                    
            slide_layout = prs.slide_layouts[1] # Title and Content
            slide = prs.slides.add_slide(slide_layout)
            
            title_shape = slide.shapes.title
            body_shape = slide.placeholders[1]
            
            if title_shape:
                title_shape.text = title
                
            if body_shape and bullets:
                tf = body_shape.text_frame
                tf.text = bullets[0]
                for bullet in bullets[1:]:
                    p = tf.add_paragraph()
                    p.text = bullet
                    p.level = 0
                    
            if notes and slide.has_notes_slide:
                notes_slide = slide.notes_slide
                notes_tf = notes_slide.notes_text_frame
                notes_tf.text = '\n'.join(notes)

        prs.save(output_path_str)
                
        if sys.stdout is None:
            import os
            os.write(1, b"OK")
        else:
            print("OK", end='')
    except Exception as e:
        msg = f"Worker 轉存 PPT 失敗: {e}\n"
        if sys.stderr is None:
            import os
            os.write(2, msg.encode('utf-8', 'replace'))
        else:
            print(msg, file=sys.stderr)
        sys.exit(1)


if __name__ == '__main__':
    if len(sys.argv) < 2:
        print('Usage: worker.py <file> [max_chars] | --ppt <output> [template]', file=sys.stderr)
        sys.exit(2)
    if sys.argv[1] == '--ppt':
        if len(sys.argv) < 3:
            print('Missing PowerPoint output path', file=sys.stderr)
            sys.exit(2)
        generate_ppt(sys.argv[2], sys.argv[3] if len(sys.argv) > 3 and sys.argv[3] else None)
    else:
        main(sys.argv[1], int(sys.argv[2]) if len(sys.argv) > 2 else 50000)
