"""Local, evidence-first document indexing for the document assistant.

The service intentionally has no network client.  It indexes only text formats
whose source locations can be preserved deterministically. Native extractors
preserve PDF pages, Word paragraphs, PowerPoint slides, and Excel ranges.
"""
from __future__ import annotations

import hashlib
import json
import logging
import re
import uuid
from datetime import datetime
from pathlib import Path

from config import settings

logger = logging.getLogger('claudecat')

DOCUMENTS_DIR = settings.LOG_DIR / 'documents'
MAX_ANALYSES_PER_DOCUMENT = 20
_SUPPORTED_SUFFIXES = {'.txt', '.md', '.csv', '.pdf', '.docx', '.pptx', '.xlsx'}
_TOKEN = re.compile(r"[A-Za-z0-9_]+|[\u4e00-\u9fff]+", re.UNICODE)
_DOCUMENT_ID_RE = re.compile(r'^\d{14}(?:_\d+)?$')
_CJK_FUNCTION_CHARS = set('的是了在有和與及或為可可以請幾多少嗎呢吧啊')


def _ensure_dir() -> None:
    DOCUMENTS_DIR.mkdir(parents=True, exist_ok=True)


def _tokens(text: str) -> set[str]:
    """Return searchable terms without treating every CJK character as proof.

    CJK text has no spaces.  Bi-grams retain useful terms such as ``付款期限``
    while avoiding a single common character (for example ``假``) matching an
    unrelated policy.
    """
    terms: set[str] = set()
    for item in _TOKEN.findall(text):
        if '\u4e00' <= item[0] <= '\u9fff':
            for index in range(len(item) - 1):
                gram = item[index:index + 2]
                if not any(char in _CJK_FUNCTION_CHARS for char in gram):
                    terms.add(gram)
        elif len(item) > 1:
            terms.add(item.lower())
    return terms


def _chunks(text: str, suffix: str) -> list[dict]:
    """Split text into small, line-addressable evidence blocks."""
    lines = text.splitlines()
    chunks: list[dict] = []
    heading = ''
    buffer: list[str] = []
    start = 1

    def flush(end: int) -> None:
        nonlocal buffer, start
        body = '\n'.join(buffer).strip()
        if body:
            chunks.append({
                'text': body,
                'source': {
                    'kind': 'line_range',
                    'locator': f'第 {start}–{end} 行' + (f' · {heading}' if heading else ''),
                    'heading': heading,
                    'line_start': start,
                    'line_end': end,
                },
            })
        buffer = []

    for line_no, line in enumerate(lines, start=1):
        if suffix == '.md' and line.startswith('#'):
            flush(line_no - 1)
            heading = line.lstrip('#').strip()
            start = line_no + 1
            continue
        if not line.strip() and buffer:
            flush(line_no - 1)
            start = line_no + 1
            continue
        if not buffer:
            start = line_no
        buffer.append(line)
        if len(buffer) >= 12:
            flush(line_no)
            start = line_no + 1
    flush(len(lines))
    return chunks


def _index_path(document_id: str) -> Path:
    return DOCUMENTS_DIR / f'{document_id}.json'


def _source(kind: str, locator: str, **values) -> dict:
    return {'kind': kind, 'locator': locator, **values}


def _extract_pdf(path: Path) -> list[dict]:
    from pypdf import PdfReader
    chunks = []
    for page_no, page in enumerate(PdfReader(path).pages, start=1):
        text = (page.extract_text() or '').strip()
        if text:
            chunks.append({'text': text, 'source': _source('pdf_page', f'第 {page_no} 頁', page=page_no)})
    if not chunks:
        raise ValueError('此掃描型 PDF 需先經 OCR 才能閱讀。')
    return chunks


def _extract_docx(path: Path) -> list[dict]:
    from docx import Document
    chunks = []
    heading = ''
    paragraph_no = 0
    for paragraph in Document(path).paragraphs:
        text = paragraph.text.strip()
        if not text:
            continue
        paragraph_no += 1
        style_name = (paragraph.style.name or '').lower() if paragraph.style else ''
        if style_name.startswith('heading'):
            heading = text
            continue
        locator = f'第 {paragraph_no} 段' + (f' · {heading}' if heading else '')
        chunks.append({'text': text, 'source': _source(
            'word_paragraph', locator, heading=heading, paragraph=paragraph_no)})
    for table_no, table in enumerate(Document(path).tables, start=1):
        rows = [[cell.text.strip() for cell in row.cells] for row in table.rows]
        headers = rows[0] if rows else []
        for row_no, cells in enumerate(rows, start=1):
            if not any(cells):
                continue
            if row_no == 1:
                text = ' | '.join(cells)
            else:
                text = ' | '.join(
                    f'{headers[index] if index < len(headers) and headers[index] else index + 1}: {value}'
                    for index, value in enumerate(cells) if value)
            chunks.append({'text': text, 'source': _source(
                'word_table', f'表格 {table_no} · 第 {row_no} 列',
                table=table_no, row=row_no, heading=heading)})
    if not chunks:
        raise ValueError('Word 文件沒有可擷取文字。')
    return chunks


def _extract_pptx(path: Path) -> list[dict]:
    from pptx import Presentation
    chunks = []
    for slide_no, slide in enumerate(Presentation(path).slides, start=1):
        text_parts = [shape.text.strip() for shape in slide.shapes if hasattr(shape, 'text') and shape.text.strip()]
        title = slide.shapes.title.text.strip() if slide.shapes.title and slide.shapes.title.text else ''
        locator = f'投影片 {slide_no}' + (f' · {title}' if title else '')
        text_extracted = bool(text_parts)
        text = '\n'.join(text_parts) if text_parts else (
            '此投影片沒有可擷取文字；圖片與圖表未分析。')
        chunks.append({'text': text, 'source': _source(
            'powerpoint_slide', locator, slide=slide_no, heading=title,
            text_extracted=text_extracted)})
    if not chunks:
        raise ValueError('PowerPoint 沒有可擷取文字。')
    return chunks


def _extract_xlsx(path: Path) -> list[dict]:
    from openpyxl import load_workbook
    workbook = load_workbook(path, read_only=True, data_only=True)
    chunks = []
    for worksheet in workbook.worksheets:
        headers: list[str] = []
        for row_no, row in enumerate(worksheet.iter_rows(values_only=True), start=1):
            cells = [(index, str(value).strip()) for index, value in enumerate(row, start=1)
                     if value is not None and str(value).strip()]
            if not cells:
                continue
            last_column = cells[-1][0]
            cell_range = f'A{row_no}:{_column_name(last_column)}{row_no}'
            locator = f'工作表 {worksheet.title} · {cell_range}'
            values = [value for _, value in cells]
            if not headers:
                headers = [str(value).strip() if value is not None else '' for value in row]
                text = ' | '.join(values)
            else:
                text = ' | '.join(
                    f'{headers[index - 1] if index <= len(headers) and headers[index - 1] else _column_name(index)}: {value}'
                    for index, value in cells)
            chunks.append({'text': text, 'source': _source(
                'excel_range', locator, sheet=worksheet.title, cell_range=cell_range)})
    workbook.close()
    if not chunks:
        raise ValueError('Excel 沒有可擷取儲存格內容。')
    return chunks


def _column_name(index: int) -> str:
    result = ''
    while index:
        index, remainder = divmod(index - 1, 26)
        result = chr(65 + remainder) + result
    return result


def _extract_chunks(path: Path, suffix: str) -> list[dict]:
    if suffix in {'.txt', '.md', '.csv'}:
        try:
            return _chunks(path.read_text(encoding='utf-8-sig'), suffix)
        except UnicodeDecodeError as exc:
            raise ValueError('此文字檔不是 UTF-8 編碼，無法安全建立本機索引。') from exc
    extractors = {'.pdf': _extract_pdf, '.docx': _extract_docx, '.pptx': _extract_pptx, '.xlsx': _extract_xlsx}
    return extractors[suffix](path)


def ingest(path_value: str) -> dict:
    path = Path(path_value)
    if not path.is_file():
        return {'error': '找不到選取的檔案。'}
    suffix = path.suffix.lower()
    if suffix not in _SUPPORTED_SUFFIXES:
        return {'error': (
            f'{suffix or "此格式"} 尚未啟用可驗證來源定位。'
            '支援 TXT、Markdown、CSV、PDF、DOCX、PPTX 與 XLSX。'
        )}
    try:
        chunks = _extract_chunks(path, suffix)
    except (OSError, ValueError, ImportError) as exc:
        return {'error': str(exc)}
    except Exception as exc:
        return {'error': f'無法解析 {suffix.upper().lstrip(".")}：{exc}'}
    if not chunks:
        return {'error': '文件沒有可建立索引的文字內容。'}

    _ensure_dir()
    digest = _file_digest(path)
    existing = _find_by_digest(digest)
    if existing is not None:
        # Same bytes: reuse the index so saved analyses stay reachable. Picking
        # the same file twice used to mint a second id, leaving two rows with
        # an identical label and the analyses attached to only one of them.
        if existing['name'] != path.name:
            existing['name'] = path.name   # renamed on disk; show the current name
            _write_index(existing)
        return {'document': _summary(existing), 'reused': True}

    document_id = _new_document_id()
    document = {
        'id': document_id,
        'name': path.name,
        'suffix': suffix,
        'digest': digest,
        'indexed_at': datetime.now().isoformat(timespec='seconds'),
        'chunks': chunks,
    }
    _write_index(document)
    return {'document': _summary(document), 'reused': False}


def _write_index(document: dict) -> None:
    _index_path(document['id']).write_text(
        json.dumps(document, ensure_ascii=False), encoding='utf-8')


def _file_digest(path: Path) -> str:
    digest = hashlib.sha256()
    with open(path, 'rb') as handle:
        for block in iter(lambda: handle.read(1 << 20), b''):
            digest.update(block)
    return digest.hexdigest()


def _find_by_digest(digest: str) -> dict | None:
    for path in DOCUMENTS_DIR.glob('*.json'):
        try:
            document = json.loads(path.read_text(encoding='utf-8'))
        except (OSError, ValueError):
            continue
        if isinstance(document, dict) and document.get('digest') == digest:
            return document
    return None


def _new_document_id() -> str:
    """YYYYMMDDHHMMSS, so the id itself says when the document was indexed.

    Hours are part of the key on purpose: minute+second alone would collide
    across the same clock minute of different hours.
    """
    base = datetime.now().strftime('%Y%m%d%H%M%S')
    candidate, ordinal = base, 0
    while _index_path(candidate).exists():   # two files indexed within one second
        ordinal += 1
        candidate = f'{base}_{ordinal}'
    return candidate


def safe_document_id(document_id: str) -> str | None:
    """Reject anything that could escape DOCUMENTS_DIR.

    Accepts both the timestamp keys minted now and the uuid4 keys that indexes
    created before 7.2.3 still carry, so existing documents keep working.
    """
    text = str(document_id)
    if _DOCUMENT_ID_RE.match(text):
        return text
    try:
        return str(uuid.UUID(text))
    except ValueError:
        return None


def _load(document_id: str) -> dict | None:
    key = safe_document_id(document_id)
    if key is None:
        return None
    try:
        return json.loads(_index_path(key).read_text(encoding='utf-8'))
    except (OSError, ValueError):
        return None


def _summary(document: dict, indexed_at: str = '') -> dict:
    return {
        'id': document['id'],
        'name': document['name'],
        'chunk_count': len(document['chunks']),
        # Shown in the list so same-named indexes stay distinguishable.
        'indexed_at': document.get('indexed_at') or indexed_at,
    }


def list_documents() -> list[dict]:
    _ensure_dir()
    result = []
    for path in DOCUMENTS_DIR.glob('*.json'):
        try:
            document = json.loads(path.read_text(encoding='utf-8'))
            # Indexes written before 7.2.3 have no indexed_at; the file's own
            # mtime is close enough to tell two same-named rows apart.
            fallback = datetime.fromtimestamp(path.stat().st_mtime).isoformat(timespec='seconds')
            result.append(_summary(document, fallback))
        except (OSError, ValueError, KeyError):
            continue
    return sorted(result, key=lambda item: (item['name'].lower(), item['indexed_at']))


def remove(document_id: str) -> dict:
    key = safe_document_id(document_id)
    if key is None:
        return {'error': '無效的文件識別碼。'}
    _index_path(key).unlink(missing_ok=True)
    # Removing the index must take the saved analyses with it, or they linger
    # as orphans no UI can reach.
    _analyses_path(key).unlink(missing_ok=True)
    return {'status': 'ok'}


def _analyses_path(document_id: str) -> Path:
    """Derived from DOCUMENTS_DIR on each call, not cached at import: tests
    repoint DOCUMENTS_DIR, and a cached value would write to the real profile.

    A subdirectory rather than "<id>.analyses.json" because list_documents()
    globs *.json here and would otherwise parse every analysis file.
    """
    return DOCUMENTS_DIR / 'analyses' / f'{document_id}.json'


def load_analyses(document_id: str) -> list[dict]:
    """Saved analyses for one document, oldest first; [] when there are none."""
    key = safe_document_id(document_id)
    if key is None:
        return []
    try:
        data = json.loads(_analyses_path(key).read_text(encoding='utf-8'))
    except (OSError, ValueError):
        return []
    return data if isinstance(data, list) else []


def latest_analysis(document_id: str) -> dict | None:
    entries = load_analyses(document_id)
    return entries[-1] if entries else None


def save_analysis(document_id: str, kind: str, label: str, result: dict) -> None:
    """Persist one analysis result so it survives closing the app.

    Stored beside the index rather than inside it: re-ingesting a document
    rewrites the index, and past analyses should outlive that.
    """
    key = safe_document_id(document_id)
    if key is None:
        return
    entries = load_analyses(key)
    entries.append({
        'kind': kind,
        'label': label,
        'saved_at': datetime.now().isoformat(timespec='seconds'),
        'answer': result.get('answer', ''),
        'sources': result.get('sources', []),
        'coverage': result.get('coverage', {}),
    })
    del entries[:-MAX_ANALYSES_PER_DOCUMENT]  # bounded: no unbounded growth on disk
    path = _analyses_path(key)
    path.parent.mkdir(parents=True, exist_ok=True)
    temporary = path.with_suffix('.tmp')
    try:
        temporary.write_text(json.dumps(entries, ensure_ascii=False), encoding='utf-8')
        temporary.replace(path)
    except OSError:
        logger.exception('could not save document analysis')
        temporary.unlink(missing_ok=True)


def query(document_id: str, question: str, limit: int = 3) -> dict:
    document = _load(document_id)
    if document is None:
        return {'error': '找不到文件索引。'}
    terms = _tokens(question)
    if not terms:
        return {'answer': '請輸入較具體的問題。', 'sources': []}

    minimum_score = 2 if len(terms) >= 2 else 1
    ranked = []
    for chunk in document['chunks']:
        score = len(terms & _tokens(chunk['text']))
        if score >= minimum_score:
            ranked.append((score, chunk))
    ranked.sort(key=lambda item: item[0], reverse=True)
    if not ranked:
        return {'answer': '此文件沒有描述此問題，無法依文件確認。', 'sources': []}

    sources = []
    for _, chunk in ranked[:limit]:
        source = dict(chunk['source'])
        source['document_name'] = document['name']
        sources.append({'excerpt': chunk['text'], 'source': source})
    # This deliberately returns evidence rather than inventing a prose answer.
    return {'answer': '找到下列與問題相關的文件內容：', 'sources': sources}


def context(document_id: str, limit: int | None = 12) -> dict:
    """Return bounded, source-bearing evidence for document-wide actions.

    A bounded context cannot be a complete reading of a long document.  Sample
    across the document instead of silently treating its first pages as the
    whole document, and expose the coverage to the caller/UI.
    """
    document = _load(document_id)
    if document is None:
        return {'error': '找不到文件索引。'}
    chunks = document['chunks']
    if limit is None or len(chunks) <= limit:
        selected = chunks
    else:
        selected = [chunks[round(index * (len(chunks) - 1) / (limit - 1))]
                    for index in range(limit)]
    sources = []
    for chunk in selected:
        source = dict(chunk['source'])
        source['document_name'] = document['name']
        sources.append({'excerpt': chunk['text'], 'source': source})
    if not sources:
        return {'error': '文件沒有可用內容。'}
    return {
        'sources': sources,
        'coverage': {
            'included_chunks': len(sources),
            'total_chunks': len(chunks),
            'complete': len(sources) == len(chunks),
            'unit': {'.pptx': '張投影片', '.pdf': '頁'}.get(
                document.get('suffix'), '個文件區塊'),
            'limited_chunks': sum(
                1 for chunk in chunks
                if chunk['source'].get('text_extracted') is False),
        },
    }
